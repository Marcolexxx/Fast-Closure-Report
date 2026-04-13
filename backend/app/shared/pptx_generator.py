import os
import sys
import logging
import tempfile
from typing import Any, Dict, List
from pathlib import Path

logger = logging.getLogger(__name__)


def _atomic_save_pptx(prs: Any, final_path: str) -> None:
    """
    PRD §7.3 Tool 11: 原子写入 PPTX 文件。
    
    写入到同目录临时文件，再 os.rename 原子替换，防止并发读到中间态文件。
    Unix 用 fcntl.flock 文件锁；Windows 无 fcntl 则直接写。
    """
    final = Path(final_path)
    tmp_fd, tmp_path = tempfile.mkstemp(dir=final.parent, suffix=".tmp.pptx")
    try:
        if sys.platform != "win32":
            import fcntl
            with open(tmp_path, "wb") as fp:
                fcntl.flock(fp, fcntl.LOCK_EX)
                prs.save(fp)
                fcntl.flock(fp, fcntl.LOCK_UN)
        else:
            # Windows: no fcntl — close the fd first, then save
            os.close(tmp_fd)
            tmp_fd = -1
            prs.save(tmp_path)
        os.rename(tmp_path, final_path)
        logger.info("pptx_atomic_write_ok", extra={"path": final_path})
    except Exception:
        # Clean up temp on failure
        try:
            if tmp_fd >= 0:
                os.close(tmp_fd)
            os.unlink(tmp_path)
        except OSError:
            pass
        raise
    finally:
        # Ensure fd closed (may already be closed)
        if tmp_fd >= 0:
            try:
                os.close(tmp_fd)
            except OSError:
                pass


def draw_bounding_boxes(image_path: str, boxes: List[Dict], out_path: str) -> bool:
    """
    Draws red bounding boxes on an image based on AI detections.
    Requires Pillow.
    boxes format: [{"box_2d": [x1, y1, x2, y2], "label": "..."}]
    """
    try:
        from PIL import Image, ImageDraw  # type: ignore
        img = Image.open(image_path)
        draw = ImageDraw.Draw(img)
        
        for box in boxes:
            coords = box.get("box_2d")
            if coords and len(coords) == 4:
                draw.rectangle(coords, outline="red", width=3)
                label = box.get("label")
                if label:
                    draw.text((coords[0], max(0, coords[1] - 15)), label, fill="red")
        
        img.save(out_path)
        return True
    except Exception as e:
        logger.error(f"Failed to draw bounding box on {image_path}: {e}")
        return False

def generate_report_pptx(task_id: str, items: List[Dict], receipts: dict, template_path: str = "") -> str:
    """
    Core PPTX renderer generating the full report and applying bounding boxes.
    Outputs to FILE_STORAGE_ROOT/aicopilot/<task_id>/
    PRD §7.3 Tool 11: uses atomic write (temp file + os.rename) to prevent dirty reads.
    """
    out_dir = Path(os.environ.get("FILE_STORAGE_ROOT", "/data")) / "aicopilot" / str(task_id)
    out_dir.mkdir(parents=True, exist_ok=True)
    
    import time
    pptx_path = str(out_dir / f"output_{int(time.time())}.pptx")
    
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
        from pptx.enum.shapes import PP_PLACEHOLDER
    except Exception as e:
        raise RuntimeError("Missing python-pptx dependency") from e

    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # Dynamic template filling logic
    # We iterate over all existing slides and shapes.
    # We replace text based on heuristics and insert pictures if we find picture/object placeholders.
    
    # Pre-compute data strings
    items_summary = []
    for it in items[:15]:
        name = it.get('name', '')
        target = it.get('target_qty', 0)
        actual = it.get('actual_qty', target)
        status = "✅ 达标" if actual >= target else "⚠️ 异常"
        items_summary.append(f"• {name} | 目标: {target} | 实际: {actual} | 状态: {status}")
    items_text = "\n".join(items_summary)
    
    matches = receipts.get("matches", [])
    unmatched = receipts.get("unmatched", [])
    fin_summary = f"匹配成功：{len(matches)} 笔 | 异常未匹配：{len(unmatched)} 笔\n"
    for m in matches[:5]:
        pmt = m.get("payment", {})
        inv = m.get("invoice", {})
        fin_summary += f"[匹配成功] 交易: {pmt.get('date')} ￥{pmt.get('amount') or inv.get('amount')} <=> 发票号: {inv.get('invoice_no')}\n"
    for u in unmatched[:5]:
        fin_summary += f"[待人工确认] 未匹配资产 | 类型: {u.get('type')} ￥{u.get('amount')}\n"

    # Available images to fill
    available_images = [it.get('design_image_path') for it in items if it.get('design_image_path') and os.path.exists(it.get('design_image_path'))]
    img_idx = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
                
            ph_type = shape.placeholder_format.type
            
            # Fill Text Placeholders (Title, Body, Center Title, Subtitle, Object, Content)
            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE):
                text = shape.text.lower()
                if "title" in text or "标题" in text or not text:
                    shape.text = "活动项目智能结案报告"
            elif ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                text = shape.text.lower()
                if "物料" in text or "清单" in text or "items" in text or "list" in text:
                    shape.text = items_text
                elif "财务" in text or "凭据" in text or "finance" in text or "receipt" in text:
                    shape.text = fin_summary
                elif not text or "content" in text:
                    shape.text = f"由 AI Copilot 自动生成 | 任务 ID: {task_id}"
            
            # Fill Picture Placeholders
            elif ph_type in (PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.OBJECT):
                if img_idx < len(available_images):
                    try:
                        shape.insert_picture(available_images[img_idx])
                        img_idx += 1
                    except Exception as e:
                        logger.error(f"Failed to insert picture into placeholder: {e}")

    # PRD §7.3 Tool 11: 原子写入，防止并发读到中间态
    _atomic_save_pptx(prs, pptx_path)
    return pptx_path

